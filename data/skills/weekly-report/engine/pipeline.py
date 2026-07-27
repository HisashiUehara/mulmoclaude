#!/usr/bin/env python3
"""WeeklyReport ゴールド取込パイプライン (engine = コード担当分).

役割分担: 抽出=コード / サニタイズ=AI / 出力=コード / 漏れスキャン=コード / 承認=人間。
このスクリプトは「抽出」「出力(在所置換+DSL)」「漏れスキャン」の3つの決定的処理を担う。
サニタイズ(架空データ差し替え)は AI が sanitized-map.json を作る別ステップ。

Subcommands:
  extract  <src>                     元ファイル(.xlsx/.pptx) -> work/extracted.json
  apply    <src> <map> <out>         原本コピーにテキストを在所置換 -> サニタイズ版 + expected.dsl.json
  scan     <src> <out>               元テキスト/数値/メタが out に残っていないか照合 -> work/leak_findings.json

書式・レイアウト・構造は「原本コピーへのテキスト在所置換」で保持し、一切再構築しない。
"""
from __future__ import annotations

import argparse
import io
import json
import re
import shutil
import sys
import zipfile
from pathlib import Path

# ---- 言語判定 (決定的) ---------------------------------------------------------
_CJK = re.compile(r"[぀-ヿ㐀-䶿一-鿿ｦ-ﾟ]")


def detect_lang(text: str) -> str:
    """CJK文字を含めば 'ja'、それ以外は 'en'。"""
    return "ja" if _CJK.search(text or "") else "en"


def pad_id(n: int) -> str:
    return f"{n:03d}"


# ---- 抽出: xlsx ---------------------------------------------------------------
def extract_xlsx(path: Path) -> dict:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=False)
    items, sheets = [], []
    idx = 0
    for ws in wb.worksheets:
        merged = [str(r) for r in ws.merged_cells.ranges]
        sheets.append({"name": ws.title, "dims": ws.dimensions, "merged": merged})
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None or str(cell.value).strip() == "":
                    continue
                f, a = cell.font, cell.alignment
                idx += 1
                items.append({
                    "id": pad_id(idx),
                    "loc": f"{ws.title}/{cell.coordinate}",
                    "text": str(cell.value),
                    "lang": detect_lang(str(cell.value)),
                    "font": {"name": f.name, "size": float(f.size) if f.size else None,
                             "bold": bool(f.bold), "italic": bool(f.italic)},
                    "align": {"h": a.horizontal, "v": a.vertical},
                    "number_format": cell.number_format,
                    "is_heading": bool(f.bold),
                })
    props = wb.properties
    meta = {k: getattr(props, k, None) for k in
            ("creator", "lastModifiedBy", "title", "subject", "description", "keywords", "category")}
    return {
        "source": {"path": str(path), "kind": "xlsx"},
        "structure": {"sheets": sheets},
        "metadata": meta,
        "items": items,
        "assets": _zip_assets(path),
    }


# ---- 抽出: pptx ---------------------------------------------------------------
def extract_pptx(path: Path) -> dict:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(path)
    items, slides = [], []
    idx = 0
    for si, slide in enumerate(prs.slides, start=1):
        shapes_meta = []
        for shi, shape in enumerate(slide.shapes, start=1):
            shapes_meta.append({"name": shape.name, "type": str(shape.shape_type)})
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    if not (run.text or "").strip():
                        continue
                    idx += 1
                    fnt = run.font
                    items.append({
                        "id": pad_id(idx),
                        "loc": f"p{si}/{shape.name}[para{pi}run{ri}]",
                        "text": run.text,
                        "lang": detect_lang(run.text),
                        "font": {"name": fnt.name,
                                 "size": fnt.size.pt if fnt.size is not None else None,
                                 "bold": bool(fnt.bold), "italic": bool(fnt.italic)},
                        "align": {"h": str(para.alignment) if para.alignment else None, "v": None},
                        "shape_pos": {"left": _emu(shape.left), "top": _emu(shape.top),
                                      "width": _emu(shape.width), "height": _emu(shape.height)},
                        "is_heading": shape.name.lower().startswith("title"),
                    })
        slides.append({"index": si, "layout": slide.slide_layout.name, "shapes": shapes_meta})
    cp = prs.core_properties
    meta = {"creator": cp.author, "lastModifiedBy": cp.last_modified_by, "title": cp.title,
            "subject": cp.subject, "keywords": cp.keywords, "category": cp.category}
    return {
        "source": {"path": str(path), "kind": "pptx"},
        "structure": {"slides": slides},
        "metadata": meta,
        "items": items,
        "assets": _zip_assets(path),
    }


def _emu(v):
    return int(v) if v is not None else None


# ---- 資産(ロゴ画像・埋込オブジェクト・ヘッダーフッター)の検出 ----------------------
def _zip_assets(path: Path) -> dict:
    images, embedded, hdrftr = [], [], []
    with zipfile.ZipFile(path) as z:
        for n in z.namelist():
            low = n.lower()
            if "/media/" in low:
                images.append(n)
            if "/embeddings/" in low or "oleobject" in low:
                embedded.append(n)
            if "header" in low or "footer" in low:
                hdrftr.append(n)
    return {"images": images, "embedded_objects": embedded, "headers_footers": hdrftr}


# ---- 抽出ディスパッチ ----------------------------------------------------------
def extract(path: Path) -> dict:
    if path.suffix.lower() == ".xlsx":
        return extract_xlsx(path)
    if path.suffix.lower() == ".pptx":
        return extract_pptx(path)
    raise SystemExit(f"unsupported file type: {path.suffix} (only .xlsx / .pptx)")


# ---- 出力: 原本コピーにテキストを在所置換 ------------------------------------------
def apply_xlsx(src: Path, out: Path, mapping: dict) -> None:
    import openpyxl

    shutil.copyfile(src, out)
    wb = openpyxl.load_workbook(out)
    for it in mapping["items"]:
        sheet, coord = it["loc"].split("/", 1)
        wb[sheet][coord] = it["sanitized_text"]  # .value のみ変更 -> 書式は保持
    p = wb.properties
    for attr in ("creator", "lastModifiedBy", "title", "subject", "description", "keywords", "category"):
        setattr(p, attr, None)
    wb.save(out)


def apply_pptx(src: Path, out: Path, mapping: dict) -> None:
    from pptx import Presentation

    shutil.copyfile(src, out)
    prs = Presentation(out)
    by_loc = {it["loc"]: it["sanitized_text"] for it in mapping["items"]}
    for si, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    loc = f"p{si}/{shape.name}[para{pi}run{ri}]"
                    if loc in by_loc:
                        run.text = by_loc[loc]  # run.font は保持される
    cp = prs.core_properties
    for attr in ("author", "last_modified_by", "title", "subject", "keywords", "category"):
        setattr(cp, attr, "")
    prs.save(out)


def apply(src: Path, out: Path, mapping: dict, keep_images: bool = False) -> list:
    if src.suffix.lower() == ".xlsx":
        apply_xlsx(src, out, mapping)
    else:
        apply_pptx(src, out, mapping)
    _scrub_zip_metadata(out)
    return _replace_media_images(out, keep_images)  # 既定: 全画像をプレースホルダへ差替え


# ---- zipレベルのメタデータ除去 (Company/Manager 等ライブラリ非公開分も) ------------
_META_TAGS = ("Company", "Manager", "dc:creator", "cp:lastModifiedBy",
              "dc:title", "dc:subject", "cp:keywords")


def _scrub_zip_metadata(path: Path) -> None:
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename in ("docProps/core.xml", "docProps/app.xml"):
                text = data.decode("utf-8", "ignore")
                for tag in _META_TAGS:
                    text = re.sub(rf"<{tag}>.*?</{tag}>", f"<{tag}></{tag}>", text, flags=re.S)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    tmp.replace(path)


# ---- 画像差し替え: 全画像を架空プレースホルダへ (デフォルト動作) ---------------------
_RASTER = ("png", "jpg", "jpeg", "gif", "bmp")
_PIL_FMT = {"png": "PNG", "jpg": "JPEG", "jpeg": "JPEG", "gif": "GIF", "bmp": "BMP"}


def _placeholder_bytes(name: str, data: bytes):
    """元画像と同サイズ・同形式のプレースホルダを生成。差替え不可(ベクタ等)は None。"""
    ext = name.rsplit(".", 1)[-1].lower()
    if ext not in _RASTER:
        return None
    try:
        from PIL import Image, ImageDraw
        w, h = Image.open(io.BytesIO(data)).size
    except Exception:
        return None
    w, h = max(1, min(w, 3000)), max(1, min(h, 3000))
    ph = Image.new("RGB", (w, h), (224, 224, 224))
    d = ImageDraw.Draw(ph)
    d.rectangle([0, 0, w - 1, h - 1], outline=(150, 150, 150))
    if w > 40 and h > 16:
        d.text((5, max(0, h // 2 - 6)), "SAMPLE (placeholder)", fill=(120, 120, 120))
    buf = io.BytesIO()
    ph.save(buf, format=_PIL_FMT[ext])
    return buf.getvalue()


def _replace_media_images(path: Path, keep_images: bool) -> list:
    """出力内の全画像(*/media/*)をプレースホルダへ差替え。keep_images=True なら保持。
    差替え不可形式(ベクタ等)は保持し 'kept-unsupported' として記録(ステップ5で目視確認)。"""
    actions = []
    tmp = path.with_suffix(path.suffix + ".imgtmp")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if "/media/" in item.filename.lower():
                if keep_images:
                    actions.append({"file": item.filename, "action": "kept"})
                else:
                    ph = _placeholder_bytes(item.filename, data)
                    if ph is not None:
                        data = ph
                        actions.append({"file": item.filename, "action": "replaced"})
                    else:
                        actions.append({"file": item.filename, "action": "kept-unsupported"})
            zout.writestr(item, data)
    tmp.replace(path)
    return actions


# ---- 出力: 対応DSL実体の組み立て (決定的) -----------------------------------------
def build_dsl(mapping: dict) -> dict:
    """sanitized-map の各項目から DSL items[] を組み立てる。
    pair_key で 原文(ja) と 訳文(en) をひとつの DSL 項目にまとめる。
    pair_key 未指定なら項目単独 (言語に応じて 原文 or 訳文 を埋める)。"""
    groups: dict[str, dict] = {}
    order: list[str] = []
    for it in mapping["items"]:
        key = it.get("pair_key") or it["id"]
        if key not in groups:
            groups[key] = {"id": key, "source_location": it["loc"], "原文": "", "訳文": "",
                           "confidence": it.get("confidence", "low")}
            order.append(key)
        g = groups[key]
        if it["lang"] == "ja":
            g["原文"] = it["sanitized_text"]
        else:
            g["訳文"] = it["sanitized_text"]
        if it.get("confidence"):
            g["confidence"] = it["confidence"]
    return {"_meta": {"generated_by": "gold-import-pipeline", "fictional": True,
                      "source_kind": mapping.get("source_kind")},
            "items": [groups[k] for k in order]}


# ---- 生XMLからのメタデータ読取り (openpyxl/pptx の既定値注入を避ける) ----------------
_META_XML_TAGS = ("dc:creator", "cp:lastModifiedBy", "dc:title", "dc:subject",
                  "cp:keywords", "Company", "Manager")


def _raw_meta(path: Path) -> dict:
    found = {}
    with zipfile.ZipFile(path) as z:
        blob = ""
        for name in ("docProps/core.xml", "docProps/app.xml"):
            if name in z.namelist():
                blob += z.read(name).decode("utf-8", "ignore")
    for tag in _META_XML_TAGS:
        m = re.search(rf"<{tag}>(.*?)</{tag}>", blob, flags=re.S)
        if m and m.group(1).strip():
            found[tag] = m.group(1).strip()
    return found


# ---- 画像カテゴリの判定 (差替え済み=ok / 保持=画像ごとに目視確認) -----------------------
def _image_finding(images_in_out: list, actions: list | None):
    if not images_in_out:
        return {"category": "ロゴ・画像", "status": "ok", "detail": "画像なし"}
    if actions:
        kept = [a for a in actions if a["action"] in ("kept", "kept-unsupported")]
        replaced = [a for a in actions if a["action"] == "replaced"]
        if not kept:
            return {"category": "ロゴ・画像", "status": "ok",
                    "detail": f"{len(replaced)}枚をプレースホルダに差替済"}
        detail = [f"目視確認: {a['file']}（"
                  + ("未対応形式・自動差替不可→手当要" if a["action"] == "kept-unsupported" else "元画像を保持")
                  + "）" for a in kept]
        return {"category": "ロゴ・画像", "status": "manual", "detail": detail}
    # actions 不明時は保守的に manual（画像ごとに目視確認）
    return {"category": "ロゴ・画像", "status": "manual",
            "detail": [f"目視確認: {img}" for img in images_in_out]}


# ---- 漏れスキャン: 実物トークンが出力に残っていないか -------------------------------
def scan(extracted: dict, out: Path, mapping: dict | None = None, images: list | None = None) -> dict:
    # 意図的に維持した汎用語(サニタイズ後==元)は漏れ扱いしない。差替え済みなのに残るものが真の漏れ。
    kept = set()
    if mapping:
        by_loc = {it["loc"]: it.get("sanitized_text", "") for it in mapping["items"]}
        for it in extracted["items"]:
            if by_loc.get(it["loc"]) == it["text"]:
                kept.add(it["text"])
    check_texts = [it["text"] for it in extracted["items"] if it["text"] not in kept]
    numbers = sorted({m for t in check_texts for m in re.findall(r"\d[\d,\.]*", t)}, key=len, reverse=True)

    out_ext = extract(out)
    out_blob = "\n".join(it["text"] for it in out_ext["items"])
    leaked_text = [t for t in check_texts if t and t in out_blob]
    leaked_nums = [n for n in numbers if len(n) >= 2 and n in out_blob]
    real_meta = _raw_meta(out)

    def cat(name, status, detail):
        return {"category": name, "status": status, "detail": detail}

    findings = [
        cat("固有名詞・テキスト", "warn" if leaked_text else "ok",
            leaked_text[:20] if leaked_text else "差替え対象テキストの残存なし"),
        cat("数値", "warn" if leaked_nums else "ok",
            leaked_nums[:20] if leaked_nums else "実物数値の残存なし"),
        cat("プロパティ情報", "warn" if real_meta else "ok",
            real_meta if real_meta else "作成者/会社/タイトル等は空(生XML照合)"),
        _image_finding(out_ext["assets"]["images"], images),
        cat("ヘッダーフッター", "manual" if out_ext["assets"]["headers_footers"] else "ok",
            out_ext["assets"]["headers_footers"] or "ヘッダーフッターなし"),
        cat("埋め込みオブジェクト", "manual" if out_ext["assets"]["embedded_objects"] else "ok",
            out_ext["assets"]["embedded_objects"] or "埋め込みオブジェクトなし"),
    ]
    ok = all(f["status"] == "ok" for f in findings)
    return {"auto_clean": ok, "kept_identical": sorted(kept), "findings": findings}


# ---- CLI ----------------------------------------------------------------------
def _load(p: str) -> dict:
    return json.loads(Path(p).read_text(encoding="utf-8"))


def _dump(obj: dict, p: Path) -> None:
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(obj, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"wrote {p}")


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="WeeklyReport gold import pipeline (code steps)")
    sub = ap.add_subparsers(dest="cmd", required=True)

    pe = sub.add_parser("extract"); pe.add_argument("src"); pe.add_argument("--out", default="engine/work/extracted.json")
    pa = sub.add_parser("apply"); pa.add_argument("src"); pa.add_argument("map"); pa.add_argument("out")
    pa.add_argument("--dsl", default="fixtures/gold/expected.dsl.json")
    pa.add_argument("--keep-images", action="store_true", help="画像を差替えず保持(ステップ5で画像ごとに目視確認)")
    pa.add_argument("--img-actions", default="engine/work/image_actions.json")
    ps = sub.add_parser("scan"); ps.add_argument("extracted"); ps.add_argument("out")
    ps.add_argument("--map", dest="map", default=None)
    ps.add_argument("--images", dest="images", default=None, help="apply が出力した image_actions.json")
    ps.add_argument("--out", dest="findings", default="engine/work/leak_findings.json")

    a = ap.parse_args(argv)
    base = Path(__file__).resolve().parent.parent  # data/skills/weekly-report

    if a.cmd == "extract":
        _dump(extract(base / a.src if not Path(a.src).is_absolute() else Path(a.src)), base / a.out)
    elif a.cmd == "apply":
        src = base / a.src if not Path(a.src).is_absolute() else Path(a.src)
        mapping = _load(a.map if Path(a.map).is_absolute() else str(base / a.map))
        mapping.setdefault("source_kind", src.suffix.lower().lstrip("."))
        out = base / a.out if not Path(a.out).is_absolute() else Path(a.out)
        out.parent.mkdir(parents=True, exist_ok=True)
        img_actions = apply(src, out, mapping, keep_images=a.keep_images)
        print(f"wrote {out}")
        _dump(build_dsl(mapping), base / a.dsl)
        _dump({"keep_images": a.keep_images, "actions": img_actions}, base / a.img_actions)
    elif a.cmd == "scan":
        ext = _load(a.extracted if Path(a.extracted).is_absolute() else str(base / a.extracted))
        out = base / a.out if not Path(a.out).is_absolute() else Path(a.out)
        mp = _load(a.map if Path(a.map).is_absolute() else str(base / a.map)) if a.map else None
        imgs = _load(a.images if Path(a.images).is_absolute() else str(base / a.images))["actions"] if a.images else None
        _dump(scan(ext, out, mp, imgs), base / a.findings)
    return 0


if __name__ == "__main__":
    sys.exit(main())
